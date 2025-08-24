from pptx import Presentation

def extract_text_from_pptx(file_path) -> tuple[list[str], list[dict]]:
    from services.parsers.text_splitter import split_text_into_chunks

    ppt = Presentation(file_path)
    all_chunks = []
    metadata = []

    for slide_num, slide in enumerate(ppt.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"

        if text.strip():
            chunks = split_text_into_chunks(text)
            all_chunks.extend(chunks)
            metadata.extend([{"slide_number": slide_num}] * len(chunks))

    return all_chunks, metadata
